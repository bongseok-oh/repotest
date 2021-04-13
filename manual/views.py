from django.core.urlresolvers import reverse_lazy
from django.db.models import Q
from django.http import HttpResponseRedirect
from django.shortcuts import redirect, render
from django.urls import reverse_lazy
from django.views.generic import CreateView, DeleteView, FormView, ListView, UpdateView

from platform import system

from konlpy.tag import Mecab # 형태소 분석을 위한 라이브러리 (Mecab 은 windows 환경에서 실행 불가)
from konlpy.tag import Kkma # For windows os
try:
    import pdftotext
except:
    pass
from pptx.dml.color import RGBColor # PPT 파일을 핸들링 하기위한 라이브러리
from pptx import Presentation
from bs4 import BeautifulSoup as bs
from tika import parser
import shutil
import os 

from .forms import *
from .models import CompareWord, WordDict

def index(request): 
    # 메인 리스트 뷰  (FBV)
    currentUser = request.session.get('adid', False)        #bongseok.oh 권한 통제를 위한 사용자 정보 가져오기 
    template_name = "manual/check.html"
    if request.method == "POST": # 파일 업로드 요청 시
        # print(request.FILES)
        try:
            newPost = CompareWord(
                tester = request.POST['tester'],
                desc = request.POST['desc'],
                origin_file = request.FILES['test_file'],
                comp = request.POST['comp'],
            )
            newPost.save() # 파일 및 정보를 저장
        except:
            return redirect(reverse_lazy('list')) # 메인 리스트로 이동(리프레시)
        return redirect(reverse_lazy('list')) # 메인 리스트로 이동(리프레시)
    else: # 페이지 로드 시
        forms = FileListForm()
        object_list = CompareWord.objects.order_by('-id') # 히스토리를 DB에서 가져옴
        return render(request, template_name, {'forms':forms, 'object_list':object_list, 'user':currentUser}) #bongseok.oh 권한 통제를 위해 currnetUser 받음

def word_list(request): # 단어 리스트 뷰 (FBV)
    currentUser = request.session.get('adid', False)
    template_name = "manual/word_list.html"
    if request.GET:
        try:

            word = request.GET['word'] # 검색 단어 
        except:
            word = ''
    else:
        word = ''
    models = WordDict.objects.all().filter(Q(word__icontains=word)).order_by('-id')  # DB에서 검색하여 저장
    context = {'object_list':models, 'word':word, 'user':currentUser}
    return render(request, template_name, context)

class WordDictCreateView(CreateView): # 단어 생성 뷰 (CBV)
    model = WordDict
    template_name = "manual/word_modify.html"
    fields = ['word','word_type', 'desc']
    success_url = reverse_lazy('word_list')

class WordDictUpdateView(UpdateView): # 단어 업데이트 뷰 (CBV)
    model = WordDict
    template_name = "manual/word_modify.html"
    fields = ['word', 'word_type', 'desc']
    success_url = reverse_lazy('word_list')

class WordDictDeleteView(DeleteView): # 단어 삭제 뷰 (CBV)
    model = WordDict
    template_name = "manual/word_confirm_delete.html"
    success_url = reverse_lazy('word_list')


def compare(request, pk): # 단어 Compare 로직
    querySet = CompareWord.objects.get(pk=pk) 
    # ref_words = [word['word'] for word in WordDict.objects.all().values('word')] # 단어 사전 리스트 Get
    ref_word = WordDict.objects.all().exclude(word_type='JTBC').values_list('word', flat=True) # 단어 사전 리스트 Get (본부 구분 전에 사용하던거)
    korean_words = WordDict.objects.all().filter(word_type='BASE').values_list('word', flat=True) # 국립국어원 단어 Get
    he_words = WordDict.objects.all().filter(word_type='HE').values_list('word', flat=True) # HE본부용 단어 Get
    ha_words = WordDict.objects.all().filter(word_type='HA').values_list('word', flat=True) # HA본부용 단어 Get
    mc_words = WordDict.objects.all().filter(word_type='MC').values_list('word', flat=True) # MC본부용 단어 Get
    lg_names = WordDict.objects.all().filter(word_type='LGE_name').values_list('word', flat=True) # LG전자 기능 고유명사 Get
    he_badwords = WordDict.objects.all().filter(word_type='BADWORD_HE').values_list('word', flat=True) # HE 지정 이슈 단어 Get
    ha_badwords = WordDict.objects.all().filter(word_type='BADWORD_HA').values_list('word', flat=True) # HA 지정 이슈 단어 Get
    mc_badwords = WordDict.objects.all().filter(word_type='BADWORD_MC').values_list('word', flat=True) # MC 지정 이슈 단어 Get
    # jtbc_words = WordDict.objects.all().filter(word_type='JTBC').values_list('word', flat=True) # JTBC 이슈 단어 Get
    he_jtbc = WordDict.objects.all().filter(word_type='JTBC_HE').values_list('word', flat=True) # JTBC_HE 이슈 단어 Get
    ha_jtbc = WordDict.objects.all().filter(word_type='JTBC_HA').values_list('word', flat=True) # JTBC_HA 이슈 단어 Get
    mc_jtbc = WordDict.objects.all().filter(word_type='JTBC_MC').values_list('word', flat=True) # JTBC_MC 이슈 단어 Get
    origin_file = querySet.origin_file # 업로드된 파일 
    result_file = 'Result_'+origin_file.name.split('/')[-1] # 저장된 파일 이름
    result_file_path = '/'.join(origin_file.name.split('/')[:-1]) # 표시 될 파일의 경로
    pptx_path = '\\'.join(origin_file.path.split('\\')[:-1]) + '\\' + result_file # DB 에 적재될 절대 경로
    paragraph_list = [] # PPT의 문단을 담기 위한 빈 List

    common_words = korean_words | lg_names | he_words | ha_words | mc_words        # 공통 선택했을 때 기본 단어셋 (국립국어원 + HE + HA + MC)
    common_he = korean_words | lg_names | he_words                                 # HE 선택 했을 때 기본 단어셋 (국립국어원 + HE)
    common_ha = korean_words | lg_names | ha_words                                 # HA 선택 했을 때 기본 단어셋 (국립국어원 + HA)
    common_mc = korean_words | lg_names | mc_words                                 # MC 선택 했을 때 기본 단어셋 (국립국어원 + MC)

    common_jtbc = he_jtbc | ha_jtbc | mc_jtbc                                      # 공통 선택했을 때 jtbc 단어셋 (HE + HA + MC)
    common_badwords = he_badwords | ha_badwords | mc_badwords                      # 공통 선택했을 때 LGE 이슈단어셋 (HE + HA + MC)

    company = [comp['comp'] for comp in CompareWord.objects.all().values('comp').order_by('-date')]     # DB에서 시간상 가장 최근 정렬로 comp 값 정렬
    company = company[0]                                                                                # comp 정렬된 것 중에 가장 최근 값 가져옴

    if str(company) == 'COMMON':
        ref_words = common_words
        jtbc_words = common_jtbc
        lge_badwords = common_badwords
    elif str(company) == 'HE':
        ref_words = common_he
        jtbc_words = he_jtbc
        lge_badwords = he_badwords
    elif str(company) == 'H&A':
        ref_words = common_ha
        jtbc_words = ha_jtbc
        lge_badwords = ha_badwords
    elif str(company) == 'MC':
        ref_words = common_mc
        jtbc_words = mc_jtbc
        lge_badwords = mc_badwords
    else:
        ref_words = ref_word
        jtbc_words = common_jtbc
        lge_badwords = common_badwords

#################################
########### pptx case ###########
#################################

    if str(origin_file).split('.')[-1] == 'pptx':
        try:
            ptx = Presentation(origin_file) # PPT 파일 로드
        except Exception as e:
            print(e)
            return redirect(reverse_lazy('list'))
        # PPT 구조
        # text frame : 슬라이드 > 쉐입 > 텍스트 프레임 > 패러그랩스 > 런 > 텍스트
        # table : 슬라이드 > 쉐입 > 테이블 > 행/열 > 패러그랩스 > 런 > 텍스트
        for slide in ptx.slides: # 슬라이드 
            for shape in slide.shapes: # 쉐입 
                if shape.has_text_frame: # Text 인경우
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            paragraph_list.append(run.text.lower())
                if shape.has_table: # table 인 경우
                    max_rows = len(shape.table.rows)
                    max_cols = len(shape.table.columns)
                    for row in range(max_rows):
                        for col in range(max_cols):
                            for paragraph in shape.table.rows[row].cells[col].text_frame.paragraphs:
                                for run in paragraph.runs:
                                    paragraph_list.append(run.text.lower())
                
        # LG 마케팅 용어 제외
        join_text = ' '.join(paragraph_list)
        join_text = join_text.lower()

        paragraphs = join_text.replace("  "," ")
        paragraphs = paragraphs.replace("  "," ")
        paragraphs = paragraphs.replace("hdmi 울트라 hd deep color","hdmi 울트라hd deep color")

        lg_name = []

        for name in lg_names:                   #LG 마케팅 용어를 모두 소문자로 변경하기 위함
            lg_name.append(name.lower())

        for _name in lg_name:                   #LG 마케팅 DB에 있는 단어면 추출된 문장에서 삭제
            paragraphs = paragraphs.replace(_name,"")
   
        if system() == 'Windows':
            word_analyzer = Kkma()
            word_index = ['NNG','NNP','OL']
        else:
            word_analyzer = Mecab()
            word_index = ['NNG','NNP','SL']

        _words = word_analyzer.pos(paragraphs) # pos 메소드 사용

        words = [] # 최종 단어를 담을 List 생성
        for _word in _words:               #분석된 단어가 두글자 이상일 때만 저장
            if len(_word[0]) > 1:
                if _word[1] in word_index: # 형태소 분석된 결과 중 word_index의 인덱스가 일치하면
                    words.append(_word[0].lower()) # 최종 단어 리스트에 추가

        ref_change = []
        out_change = []

        for ref in ref_words:                   #ref_words의 단어를 모두 소문자로 변경하기 위함
            ref_change.append(ref.lower())
        
        ref_words = ref_change

        for out in words:                       #pdf는 텍스트 추출시 대소문자 구분되어 소문자로 변경해야함
            out_change.append(out.lower())
        
        words = out_change

        bad_words = set(words) - set(ref_words) # 단어 사전과 PPT에서 검색된 단어와 차집합을 구함 > BAD Words
        print('bad words count : {}'.format(len(bad_words)))

        # JTBC/LGE 이슈 단어셋을 소문자로 변경
        jtbc_word = []
        lge_badword = []

        for j in jtbc_words:
            jtbc_word.append(j.lower())
        
        for l in lge_badwords:
            lge_badword.append(l.lower())
            
        # JTBC Bad words와 LGE 이슈단어, 국립국어원 Bad_words를 구분해서 저장
        bad_words_jtbc = []
        bad_words_lge = []

        set_words = set(words)
        list_words = list(set_words)
        
        for check in list_words:
            if check in jtbc_word:
                bad_words_jtbc.append(check)                            # JTBC Words만 따로 구성
            elif check in lge_badword:
                bad_words_lge.append(check)                             # LGE 이슈 단어만 따로 구성

        bad_words_kor = set(bad_words) - set(bad_words_jtbc) - set(bad_words_lge)           # JTBC 제외한 나머지 Bad Words만 따로 구성

        bad_words_total = set(bad_words_kor) | set(bad_words_jtbc) | set(bad_words_lge)     # Total Bad Words 구성 (최종 Bad Words 단어수 등에서 사용)

       # JTBC 단어와 ref_DB 단어 각각 카운트 해주기 위해 
        _bad_total_kor = 0
        _bad_total_jtbc = 0
        _bad_total_lge = 0
        _bad_words_kor = ''
        _bad_words_jtbc = ''
        _bad_words_lge = ''

        for b in bad_words_kor:
            cnt=words.count(b)
            _bad_total_kor += cnt
            _bad_words_kor += '{}({}), '.format(b, cnt)
 
        for c in bad_words_jtbc:
            cnt=words.count(c)
            _bad_total_jtbc += cnt
            _bad_words_jtbc += '{}({}), '.format(c, cnt)

        for d in bad_words_lge:
            cnt=words.count(d)
            _bad_total_lge += cnt
            _bad_words_lge += '{}({}), '.format(d, cnt)

        _bad_total = _bad_total_kor + _bad_total_jtbc + _bad_total_lge

        results = ''

        if len(bad_words_total) == 0: # 판정 결과 
            results = '<span style="font-size:1.0em; color:green;"><b>PASS</b></span>'
        else:
            results = '<span style="font-size:1.0em; color:red;"><b>NG</b></span></br>' # NG 이면 해당 단어를 찾아 PPT 파일을 수정
            results += '<span style="font-size:1.0em;"><b>Bad Words : {} (Total Words: {} / Bad Words Sum: {})</b></span></br></br>'.format(len(bad_words_total), len(words),_bad_total)
            for slide in ptx.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for word in bad_words_total:
                                    if word.lower() in run.text.lower():
                                        font = run.font
                                        font.bold = True # 굵게 
                                        font.italic = True # 기울임체
                                        font.color.rgb = RGBColor(0xFF, 0x7F, 0x50) # 주황색 
                                    
                    if shape.has_table:
                        max_rows = len(shape.table.rows)
                        max_cols = len(shape.table.columns)
                        for row in range(max_rows):
                            for col in range(max_cols):
                                for paragraph in shape.table.rows[row].cells[col].text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        for word in bad_words_total:
                                            if word.lower() in run.text.lower():
                                                font = run.font
                                                font.bold = True
                                                font.italic = True
                                                font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

        results_jtbc = ''                                                                                              #JTBC 단어는 색상 구분
        results_jtbc += '<span style="font-size:1.0em; color:#EC7513;"> {} </span></br>'.format(_bad_words_jtbc)       #JTBC 단어는 색상 구분

        results_lge = ''                                                                                               #LGE 이슈 단어는 색상 구분
        results_lge += '<span style="font-size:1.0em; color:#660033;"> {} </span></br>'.format(_bad_words_lge)         #LGE 이슈 단어는 색상 구분        
               
        result = str(results) + str(results_jtbc) + str(results_lge) + _bad_words_kor                                  #JTBC/LGE 이슈 단어 구분하여 결과 표시
   
        print(pptx_path)
        ptx.save(result_file) # PPT 파일 저장

        try:
            shutil.move(result_file, 'static/files/'+result_file_path)
        except:
            os.remove('static/files/'+result_file_path + '/' + result_file)
            shutil.move(result_file, 'static/files/'+result_file_path)
        querySet.result = result
        querySet.result_file = result_file_path+'/'+result_file
        querySet.save() # DB 저장

        return redirect(reverse_lazy('list'))

#################################
########### html case ###########
#################################

    elif str(origin_file).split('.')[-1] == 'html' or str(origin_file).split('.')[-1] == 'htm':
        decoded_html = origin_file.read().decode('utf-8')
        soup = bs(decoded_html, 'html.parser')

        content = soup.find_all(text=True)
        content = str(content)

        #LG 마케팅 용어 제외
        content = content.replace("\\ufeff"," ")
        content = content.replace(r"\r\n",r" ")
        contents = content.lower() 

        lg_name = []

        for name in lg_names:                   #LG 마케팅 용어를 모두 소문자로 변경하기 위함
            lg_name.append(name.lower()) 

        for _name in lg_name:
            contents = contents.replace(_name,"") 

        if system() == 'Windows':
            word_analyzer = Kkma()
            word_index = ['NNG','NNP','OL']
        else:
            word_analyzer = Mecab()
            word_index = ['NNG','NNP','SL']

        _words = word_analyzer.pos(contents) # pos 메소드 사용

        words = [] # 최종 단어를 담을 List 생성
        for _word in _words:               #분석된 단어가 두글자 이상일 때만 저장
            if len(_word[0]) > 1:
                if _word[1] in word_index: # 형태소 분석된 결과 중 word_index의 인덱스가 일치하면
                    words.append(_word[0].lower()) # 최종 단어 리스트에 추가(소문자로 해서  by David)

        ref_change = []
        out_change = []

        for ref in ref_words:                   #ref_words의 단어를 모두 소문자로 변경하기 위함
            ref_change.append(ref.lower())
        
        ref_words = ref_change

        for out in words:                       
            out_change.append(out.lower())
        
        words = out_change

        bad_words = set(words) - set(ref_words)

        # JTBC/LGE 이슈 단어셋을 소문자로 변경
        jtbc_word = []
        lge_badword = []

        for j in jtbc_words:
            jtbc_word.append(j.lower())
        
        for l in lge_badwords:
            lge_badword.append(l.lower())
            
        # JTBC Bad words와 LGE 이슈단어, 국립국어원 Bad_words를 구분해서 저장
        bad_words_jtbc = []
        bad_words_lge = []

        set_words = set(words)
        list_words = list(set_words)
        
        for check in list_words:
            if check in jtbc_word:
                bad_words_jtbc.append(check)                            # JTBC Words만 따로 구성
            elif check in lge_badword:
                bad_words_lge.append(check)                             # LGE 이슈 단어만 따로 구성

        bad_words_kor = set(bad_words) - set(bad_words_jtbc) - set(bad_words_lge)           # JTBC 제외한 나머지 Bad Words만 따로 구성

        bad_words_total = set(bad_words_kor) | set(bad_words_jtbc) | set(bad_words_lge)     # Total Bad Words 구성 (최종 Bad Words 단어수 등에서 사용)

       # JTBC 단어와 ref_DB 단어 각각 카운트 해주기 위해 
        _bad_total_kor = 0
        _bad_total_jtbc = 0
        _bad_total_lge = 0
        _bad_words_kor = ''
        _bad_words_jtbc = ''
        _bad_words_lge = ''

        for b in bad_words_kor:
            cnt=words.count(b)
            _bad_total_kor += cnt
            _bad_words_kor += '{}({}), '.format(b, cnt)
 
        for c in bad_words_jtbc:
            cnt=words.count(c)
            _bad_total_jtbc += cnt
            _bad_words_jtbc += '{}({}), '.format(c, cnt)

        for d in bad_words_lge:
            cnt=words.count(d)
            _bad_total_lge += cnt
            _bad_words_lge += '{}({}), '.format(d, cnt)

        _bad_total = _bad_total_kor + _bad_total_jtbc + _bad_total_lge

        results = ''

        if len(bad_words_total) == 0: 
            results = '<span style="font-size:1.0em; color:green;"><b>PASS</b></span>'
        else:
            results = '<span style="font-size:1.0em; color:red;"><b>NG</b></span></br>'
            results += '<span style="font-size:1.0em;"><b>Bad Words : {} (Total Words: {} / Bad Words Sum: {})</b></span></br></br>'.format(len(bad_words_total), len(words),_bad_total)
        
        results_jtbc = ''                                                                                              #JTBC 단어는 색상 구분
        results_jtbc += '<span style="font-size:1.0em; color:#EC7513;"> {} </span></br>'.format(_bad_words_jtbc)       #JTBC 단어는 색상 구분

        results_lge = ''                                                                                               #LGE 이슈 단어는 색상 구분
        results_lge += '<span style="font-size:1.0em; color:#660033;"> {} </span></br>'.format(_bad_words_lge)         #LGE 이슈 단어는 색상 구분        
               
        result = str(results) + str(results_jtbc) + str(results_lge) + _bad_words_kor                                  #JTBC/LGE 이슈 단어 구분하여 결과 표시
 
        html_result = decoded_html
        _soup = bs(html_result, 'html.parser')
        _contents= _soup.get_text()
        
        for bad_word in bad_words_total:
            bad_word = bad_word.lower()                  
            _contents = _contents.replace(bad_word, '<b><font style="color:rgb(255,153,0);">{}</font></b>'.format(bad_word))

        html_result_file = open(result_file, 'w')
        html_result_file.write(_contents)
        html_result_file.close()

        try:
            shutil.move(result_file, 'static/files/'+result_file_path)
        except:
            os.remove('static/files/'+result_file_path + '/' + result_file)
            shutil.move(result_file, 'static/files/'+result_file_path)
        querySet.result = result
        querySet.result_file = result_file_path+'/'+result_file
        querySet.save() # DB 저장

        return redirect(reverse_lazy('list'))

##################################
############ pdf case ############
##################################

    elif str(origin_file).split('.')[-1] == 'pdf':

        with open(origin_file.path, 'rb') as f:
            pdf = pdftotext.PDF(f)

        paragraph = '\n'.join(pdf)
        #LG 마케팅 용어 제외
        paragraph = paragraph.replace("\n"," ")
        paragraphs = paragraph.lower()                   #pdf는 영문이 대문자로 추출되니 소문자로 변경 필요
  
        lg_name = []

        for name in lg_names:                   #LG 마케팅 용어를 모두 소문자로 변경하기 위함
            lg_name.append(name.lower()) 

        for _name in lg_name:
            paragraphs = paragraphs.replace(_name,"") 

        if system() == 'Windows':
            word_analyzer = Kkma()
            word_index = ['NNG','NNP','OL']
        else:
            word_analyzer = Mecab()
            word_index = ['NNG','NNP','SL']

        _words = word_analyzer.pos(paragraphs) # pos 메소드 사용
        
        words = [] # 최종 단어를 담을 List 생성
        for _word in _words:               #분석된 단어가 두글자 이상일 때만 저장
            if len(_word[0]) > 1:
                if _word[1] in word_index: # 형태소 분석된 결과 중 word_index의 인덱스가 일치하면
                    words.append(_word[0].lower()) # 최종 단어 리스트에 추가

        ref_change = []
        out_change = []

        for ref in ref_words:                   #ref_words의 단어를 모두 소문자로 변경하기 위함
            ref_change.append(ref.lower())
        
        ref_words = ref_change

        for out in words:                       #pdf는 텍스트 추출시 대소문자 구분되어 소문자로 변경해야함
            out_change.append(out.lower())
        
        words = out_change

        bad_words = set(words) - set(ref_words)
        print('bad words count : {}'.format(len(bad_words)))
  
        # JTBC/LGE 이슈 단어셋을 소문자로 변경
        jtbc_word = []
        lge_badword = []

        for j in jtbc_words:
            jtbc_word.append(j.lower())
        
        for l in lge_badwords:
            lge_badword.append(l.lower())
            
        # JTBC Bad words와 LGE 이슈단어, 국립국어원 Bad_words를 구분해서 저장
        bad_words_jtbc = []
        bad_words_lge = []

        set_words = set(words)
        list_words = list(set_words)
        
        for check in list_words:
            if check in jtbc_word:
                bad_words_jtbc.append(check)                            # JTBC Words만 따로 구성
            elif check in lge_badword:
                bad_words_lge.append(check)                             # LGE 이슈 단어만 따로 구성

        bad_words_kor = set(bad_words) - set(bad_words_jtbc) - set(bad_words_lge)           # JTBC 제외한 나머지 Bad Words만 따로 구성

        bad_words_total = set(bad_words_kor) | set(bad_words_jtbc) | set(bad_words_lge)     # Total Bad Words 구성 (최종 Bad Words 단어수 등에서 사용)

       # JTBC 단어와 ref_DB 단어 각각 카운트 해주기 위해 
        _bad_total_kor = 0
        _bad_total_jtbc = 0
        _bad_total_lge = 0
        _bad_words_kor = ''
        _bad_words_jtbc = ''
        _bad_words_lge = ''

        for b in bad_words_kor:
            cnt=words.count(b)
            _bad_total_kor += cnt
            _bad_words_kor += '{}({}), '.format(b, cnt)
 
        for c in bad_words_jtbc:
            cnt=words.count(c)
            _bad_total_jtbc += cnt
            _bad_words_jtbc += '{}({}), '.format(c, cnt)

        for d in bad_words_lge:
            cnt=words.count(d)
            _bad_total_lge += cnt
            _bad_words_lge += '{}({}), '.format(d, cnt)

        _bad_total = _bad_total_kor + _bad_total_jtbc + _bad_total_lge

        results = ''

        if len(bad_words_total) == 0: # 판정 결과 
            results = '<span style="font-size:1.0em; color:green;"><b>PASS</b></span>'
        else:
            results = '<span style="font-size:1.0em; color:red;"><b>NG</b></span></br>'
            results += '<span style="font-size:1.0em;"><b>Bad Words : {} (Total Words: {} / Bad Words Sum: {})</b></span></br></br>'.format(len(bad_words_total), len(words),_bad_total)

        results_jtbc = ''                                                                                              #JTBC 단어는 색상 구분
        results_jtbc += '<span style="font-size:1.0em; color:#EC7513;"> {} </span></br>'.format(_bad_words_jtbc)       #JTBC 단어는 색상 구분

        results_lge = ''                                                                                               #LGE 이슈 단어는 색상 구분
        results_lge += '<span style="font-size:1.0em; color:#660033;"> {} </span></br>'.format(_bad_words_lge)         #LGE 이슈 단어는 색상 구분        
               
        result = str(results) + str(results_jtbc) + str(results_lge) + _bad_words_kor                                  #JTBC/LGE 이슈 단어 구분하여 결과 표시
        
        querySet.result = result
        querySet.save() # DB 저장

        return redirect(reverse_lazy('list'))

def comparehtml(request, pk):

    return 'ok'


